"""文件解析入口与解析结果结构。"""

from __future__ import annotations

from dataclasses import dataclass, field
from html.parser import HTMLParser
from pathlib import Path


class _TextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        text = data.strip()
        if text:
            self.parts.append(text)


@dataclass
class ParsedDocument:
    """解析结果：全文、分页文本与标题线索。"""

    text: str
    pages: list[str] = field(default_factory=list)
    error: str | None = None


def _read_text(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gbk"):
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return path.read_text(errors="ignore")


def parse_file(path: Path) -> ParsedDocument:
    """按扩展名解析文件为纯文本。"""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return _parse_pdf(path)
        if ext == ".docx":
            return _parse_docx(path)
        if ext == ".pptx":
            return _parse_pptx(path)
        if ext in (".xlsx", ".xlsm"):
            return _parse_xlsx(path)
        if ext == ".html" or ext == ".htm":
            return _parse_html(path)
        return ParsedDocument(text=_read_text(path))
    except Exception as exc:  # noqa: BLE001 - 解析失败返回错误信息
        return ParsedDocument(text="", error=f"{type(exc).__name__}: {exc}")


def _parse_pdf(path: Path) -> ParsedDocument:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = [(page.extract_text() or "") for page in reader.pages]
    return ParsedDocument(text="\n".join(pages), pages=pages)


def _parse_docx(path: Path) -> ParsedDocument:
    from docx import Document

    document = Document(str(path))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    return ParsedDocument(text="\n".join(paragraphs))


def _parse_pptx(path: Path) -> ParsedDocument:
    from pptx import Presentation

    presentation = Presentation(str(path))
    pages: list[str] = []
    for slide in presentation.slides:
        texts = [
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        pages.append("\n".join(texts))
    return ParsedDocument(text="\n\n".join(pages), pages=pages)


def _parse_xlsx(path: Path) -> ParsedDocument:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"【工作表：{sheet.title}】")
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell) for cell in row if cell is not None]
            if values:
                lines.append(" | ".join(values))
    workbook.close()
    return ParsedDocument(text="\n".join(lines))


def _parse_html(path: Path) -> ParsedDocument:
    extractor = _TextExtractor()
    extractor.feed(_read_text(path))
    return ParsedDocument(text="\n".join(extractor.parts))
